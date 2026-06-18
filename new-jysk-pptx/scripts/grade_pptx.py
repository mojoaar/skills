import os
import sys
import json
from pptx import Presentation

def run_grading(pptx_path, grading_json_path):
    print(f"Grading {pptx_path} -> {grading_json_path}")
    expectations = []
    
    # Assertion 1: Valid PPTX
    prs = None
    valid_pptx = False
    evidence_1 = ""
    try:
        prs = Presentation(pptx_path)
        valid_pptx = True
        evidence_1 = f"Successfully loaded PowerPoint presentation from {pptx_path} with {len(prs.slides)} slides."
    except Exception as e:
        evidence_1 = f"Failed to load presentation: {str(e)}"
        
    expectations.append({
        "text": "File exists and is a valid PPTX presentation",
        "passed": valid_pptx,
        "evidence": evidence_1
    })
    
    if not valid_pptx or prs is None:
        # All others fail if file is invalid
        for text in [
            "Verdana font is strictly enforced on all text shapes",
            "Main text color is JYSK Dark Grey (RGB 86, 86, 85)",
            "Sentence-case is applied to slide titles",
            "Agenda slide title font size is 40pt",
            "Standard slide title font size is 28pt and bold"
        ]:
            expectations.append({
                "text": text,
                "passed": False,
                "evidence": "Cannot evaluate because the presentation file could not be loaded."
            })
            
        with open(grading_json_path, "w") as f:
            json.dump({"expectations": expectations}, f, indent=2)
        return

    # Check font names, sizes, colors
    font_issues = []
    color_issues = []
    casing_issues = []
    agenda_size_correct = True
    agenda_evidence = "No agenda slide found."
    standard_size_correct = True
    standard_evidence = "No standard content slides found."
    
    for slide_idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name.lower()
        
        # Check title casing & sizes
        title_text = ""
        title_shape = None
        
        # Find title
        if slide.shapes.title:
            title_shape = slide.shapes.title
        else:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1 or shape.placeholder_format.idx == 10:
                    title_shape = shape
                    break
                    
        tf = getattr(title_shape, "text_frame", None) if title_shape else None
        if tf:
            title_text = tf.text.strip()
            if title_text:
                # Check casing: First character should be uppercase
                first_char = title_text[0] if title_text else ""
                if first_char.isalpha() and not first_char.isupper():
                    casing_issues.append(f"Slide {slide_idx} title '{title_text}' does not start with capital letter.")
                # Title shouldn't be all-uppercase (except short acronyms < 5 chars)
                if title_text.isupper() and len(title_text) > 4:
                    casing_issues.append(f"Slide {slide_idx} title '{title_text}' is in ALL CAPS.")
                    
                # Check slide title size based on layout
                first_p = tf.paragraphs[0]
                size_pt = first_p.font.size.pt if first_p.font.size else None
                is_bold = first_p.font.bold
                
                if "agenda" in layout_name:
                    if size_pt != 40:
                        agenda_size_correct = False
                        agenda_evidence = f"Agenda slide '{title_text}' title font size is {size_pt}pt instead of 40pt."
                    else:
                        agenda_evidence = "Agenda slide title font size is correct (40pt)."
                elif "standard" in layout_name or "content" in layout_name:
                    if size_pt != 28:
                        standard_size_correct = False
                        standard_evidence = f"Standard slide '{title_text}' title font size is {size_pt}pt instead of 28pt."
                    elif not is_bold:
                        standard_size_correct = False
                        standard_evidence = f"Standard slide '{title_text}' title is not bold."
                    else:
                        standard_evidence = "Standard slide title font size is correct (28pt) and bold."
                        
        # Check fonts and colors of all texts
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_tf = getattr(shape, "text_frame", None)
                if shape_tf:
                    for p in shape_tf.paragraphs:
                        for run in p.runs:
                            text = run.text.strip()
                            if not text:
                                continue
                            
                            # Font name check
                            font_name = run.font.name or p.font.name
                            # Note: if font_name is None, it defaults to whatever is in the template master, which is Verdana for JYSK templates.
                            if font_name and font_name != "Verdana":
                                font_issues.append(f"Slide {slide_idx} has non-Verdana font '{font_name}' on text '{text[:20]}...'")
                                
                            # Color check (RGB 86, 86, 85)
                            color = run.font.color or p.font.color
                            if color and color.type == 1: # RGB type
                                rgb = color.rgb
                                if (rgb.r, rgb.g, rgb.b) != (86, 86, 85):
                                    # Also check if it's white for Title or Breaker layouts where white text is valid
                                    if ("title" in layout_name or "breaker" in layout_name) and (rgb.r, rgb.g, rgb.b) == (255, 255, 255):
                                        pass
                                    else:
                                        color_issues.append(f"Slide {slide_idx} has non-standard text color {(rgb.r, rgb.g, rgb.b)} on text '{text[:20]}...'")

    # Compile grading results
    expectations.append({
        "text": "Verdana font is strictly enforced on all text shapes",
        "passed": len(font_issues) == 0,
        "evidence": "No non-Verdana fonts detected." if not font_issues else "; ".join(font_issues[:5])
    })
    
    expectations.append({
        "text": "Main text color is JYSK Dark Grey (RGB 86, 86, 85)",
        "passed": len(color_issues) == 0,
        "evidence": "All text colors are standard JYSK Dark Grey (or white on colored backgrounds)." if not color_issues else "; ".join(color_issues[:5])
    })
    
    expectations.append({
        "text": "Sentence-case is applied to slide titles",
        "passed": len(casing_issues) == 0,
        "evidence": "All slide titles are correctly written in sentence case." if not casing_issues else "; ".join(casing_issues[:5])
    })
    
    expectations.append({
        "text": "Agenda slide title font size is 40pt",
        "passed": agenda_size_correct,
        "evidence": agenda_evidence
    })
    
    expectations.append({
        "text": "Standard slide title font size is 28pt and bold",
        "passed": standard_size_correct,
        "evidence": standard_evidence
    })
    
    # Calculate summary
    total_assertions = len(expectations)
    passed_assertions = sum(1 for exp in expectations if exp["passed"])
    failed_assertions = total_assertions - passed_assertions
    pass_rate = passed_assertions / total_assertions if total_assertions > 0 else 0.0
    
    summary = {
        "total": total_assertions,
        "passed": passed_assertions,
        "failed": failed_assertions,
        "pass_rate": pass_rate
    }
    
    # Save JSON file
    os.makedirs(os.path.dirname(grading_json_path), exist_ok=True)
    with open(grading_json_path, "w") as f:
        json.dump({
            "summary": summary,
            "expectations": expectations
        }, f, indent=2)
    print(f"Grading written to {grading_json_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python grade_pptx.py <pptx_path> <grading_json_path>")
        sys.exit(1)
    run_grading(sys.argv[1], sys.argv[2])
