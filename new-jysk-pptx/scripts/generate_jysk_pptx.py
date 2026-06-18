import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# JYSK Standard Formatting Constants
FONT_NAME = "Verdana"
TEXT_COLOR = (86, 86, 85)  # RGB for Hex #565655 (Dark Grey)

def to_sentence_case(text):
    """Converts text to sentence case, preserving acronyms and mixed-case words."""
    if not text:
        return text
    words = text.split()
    if not words:
        return text
    
    first_word = words[0]
    # If the first word is already all-uppercase and longer than 1, preserve it as an acronym.
    if not (first_word.isupper() and len(first_word) > 1):
        # Otherwise, capitalise the first letter and keep the rest
        first_word = first_word[0].upper() + first_word[1:] if len(first_word) > 1 else first_word.upper()
    
    # For subsequent words, preserve them exactly as they are to avoid destroying acronyms (e.g. IT, CISO) or brands (WELLPUR)
    return " ".join([first_word] + words[1:])

def find_layout_by_patterns(prs, patterns):
    """Searches for a slide layout matching any of the pattern strings."""
    for pattern in patterns:
        for layout in prs.slide_layouts:
            if pattern.lower() in layout.name.lower():
                return layout
    # Return first layout as safe fallback
    return prs.slide_layouts[0]

def format_title_shape(title_shape, text, size_pt=28):
    """Formats a title shape with the JYSK corporate font style."""
    if not title_shape:
        return
    title_shape.text = to_sentence_case(text)
    for p in title_shape.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(size_pt)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*TEXT_COLOR)

def format_paragraph(p, text, level, base_size=20):
    """Formats a paragraph with Verdana, JYSK Dark Grey, and size adjusted for level."""
    p.text = text
    p.level = level
    p.font.name = FONT_NAME
    # JYSK rule: Bullet indent reduced by 2pt for each indent, until 5th level (don't go below 12pt)
    size = max(12, base_size - (level * 2))
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*TEXT_COLOR)

def add_title_slide(prs, title, subtitle, template_type):
    """Adds a JYSK-branded Title Slide."""
    patterns = ["title page", "7_title page", "title"]
    layout = find_layout_by_patterns(prs, patterns)
    slide = prs.slides.add_slide(layout)
    
    # Title
    title_shape = None
    if slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1: # TITLE
                title_shape = shape
                break
                
    if title_shape:
        format_title_shape(title_shape, title, size_pt=28)
        
    # Subtitle
    subtitle_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 4: # SUBTITLE
            subtitle_shape = shape
            break
            
    if subtitle_shape and subtitle:
        subtitle_shape.text = subtitle
        for p in subtitle_shape.text_frame.paragraphs:
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*TEXT_COLOR)
            
    return slide

def add_agenda_slide(prs, title, items):
    """Adds a JYSK-branded Agenda Slide."""
    patterns = ["agenda", "1_agenda", "agenda "]
    layout = find_layout_by_patterns(prs, patterns)
    slide = prs.slides.add_slide(layout)
    
    # Title (size 40 per JYSK guidelines, usually idx=10 in JYSK Agenda layout)
    title_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 10:
            title_shape = shape
            break
    if not title_shape and slide.shapes.title:
        title_shape = slide.shapes.title
    if not title_shape:
        # Fallback to any BODY placeholder that isn't the main list
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [1, 2] and shape.placeholder_format.idx != 2:
                title_shape = shape
                break
                
    if title_shape:
        format_title_shape(title_shape, title or "Agenda", size_pt=40)
        
    # Content (usually idx=2 in JYSK Agenda layout)
    content_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 2:
            content_shape = shape
            break
    if not content_shape:
        # Fallback to any OBJECT or BODY placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [2, 7] and shape != title_shape:
                content_shape = shape
                break
            
    if content_shape and items:
        tf = content_shape.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            format_paragraph(p, item, level=0, base_size=20)
            
    return slide

def add_standard_slide(prs, title, body_items):
    """Adds a standard JYSK slide with Title and Bullets (supporting indents)."""
    patterns = ["standard slide", "standard"]
    layout = find_layout_by_patterns(prs, patterns)
    slide = prs.slides.add_slide(layout)
    
    # Title (size 28 bold)
    title_shape = None
    if slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1: # TITLE
                title_shape = shape
                break
                
    if title_shape:
        format_title_shape(title_shape, title, size_pt=28)
        
    # Content
    content_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type in [2, 7]: # BODY, OBJECT
            content_shape = shape
            break
            
    if content_shape and body_items:
        tf = content_shape.text_frame
        tf.clear()
        for i, item in enumerate(body_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            
            # Determine bullet level based on leading spaces (2 spaces = level 1, etc.)
            level = 0
            cleaned_item = item
            if item.startswith("  "):
                stripped = item.lstrip()
                indent_spaces = len(item) - len(stripped)
                level = min(5, indent_spaces // 2)
                cleaned_item = stripped
                
            # Remove leading bullet symbols if present since PPTX adds them automatically
            if cleaned_item.startswith("- "):
                cleaned_item = cleaned_item[2:]
            elif cleaned_item.startswith("* "):
                cleaned_item = cleaned_item[2:]
                
            format_paragraph(p, cleaned_item, level=level, base_size=20)
            
    return slide

def add_breaker_slide(prs, title, subtitle):
    """Adds a JYSK Section Breaker / Divider Slide."""
    patterns = ["breaker large placeholder white", "breaker large placeholder blue", "1_breaker large placeholder", "breaker"]
    layout = find_layout_by_patterns(prs, patterns)
    slide = prs.slides.add_slide(layout)
    
    # Title (size 28 bold)
    title_shape = None
    if slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1: # TITLE
                title_shape = shape
                break
                
    if title_shape:
        format_title_shape(title_shape, title, size_pt=28)
        
    # Subtitle
    subtitle_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 4: # SUBTITLE
            subtitle_shape = shape
            break
            
    if subtitle_shape and subtitle:
        subtitle_shape.text = subtitle
        for p in subtitle_shape.text_frame.paragraphs:
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*TEXT_COLOR)
            
    return slide

def generate_presentation(data):
    """Generates the PowerPoint presentation from structured input data."""
    template_type = data.get("template", "indoor").lower()
    
    # Determine base template path
    assets_dir = "/Users/mojoaar/AI/skills/new-jysk-pptx/assets"
    if template_type == "outdoor" or "outdoor" in template_type:
        template_path = f"{assets_dir}/jysk_outdoor.pptx"
    else:
        template_path = f"{assets_dir}/jysk_indoor.pptx"
        
    if not os.path.exists(template_path):
        # Fallback to check if any pptx is present
        print(f"Warning: Template {template_path} not found. Checking directory contents...", file=sys.stderr)
        pptx_files = [f for f in os.listdir(assets_dir) if f.endswith(".pptx")]
        if pptx_files:
            template_path = os.path.join(assets_dir, pptx_files[0])
            print(f"Using alternative template: {template_path}", file=sys.stderr)
        else:
            raise FileNotFoundError(f"No PowerPoint templates found in {assets_dir}")
            
    prs = Presentation(template_path)
    
    # Remove default slides in the template so we start fresh
    # Safe loop to delete existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
        
    # Populate slides
    for idx, slide_data in enumerate(data.get("slides", [])):
        stype = slide_data.get("type", "standard").lower()
        title = slide_data.get("title", f"Slide {idx + 1}")
        
        if stype == "title":
            add_title_slide(prs, title, slide_data.get("subtitle", ""), template_type)
        elif stype == "agenda":
            add_agenda_slide(prs, title, slide_data.get("items", []))
        elif stype == "breaker" or stype == "divider":
            add_breaker_slide(prs, title, slide_data.get("subtitle", ""))
        else: # Standard Slide
            body = slide_data.get("body", [])
            # Also handle if a list is provided as a newline separated string
            if isinstance(body, str):
                body = body.split("\n")
            add_standard_slide(prs, title, body)
            
    output_path = data.get("output_path", "jysk_presentation.pptx")
    prs.save(output_path)
    print(f"Success: Presentation saved to {output_path}")

def parse_args():
    parser = argparse.ArgumentParser(description="JYSK PowerPoint Generator")
    parser.add_argument("--input", help="Path to input JSON file containing slide structure")
    parser.add_argument("--template", default="indoor", choices=["indoor", "outdoor"], help="JYSK template type")
    parser.add_argument("--output", default="jysk_presentation.pptx", help="Path to output PPTX file")
    
    # Fast CLI arguments
    parser.add_argument("--title", help="Fast Title slide creation")
    parser.add_argument("--subtitle", default="", help="Fast Subtitle")
    parser.add_argument("--slides", help="Fast slides formatted as 'Title1|B1|B2;Title2|B1'")
    
    return parser.parse_args()

def main():
    args = parse_args()
    
    if args.input:
        # Load from JSON file
        with open(args.input, "r") as f:
            data = json.load(f)
    elif args.title or args.slides:
        # Construct slide data from fast CLI arguments
        data = {
            "template": args.template,
            "output_path": args.output,
            "slides": []
        }
        if args.title:
            data["slides"].append({
                "type": "title",
                "title": args.title,
                "subtitle": args.subtitle
            })
            
        if args.slides:
            slide_groups = args.slides.split(";")
            for group in slide_groups:
                parts = group.split("|")
                slide_title = parts[0]
                bullets = parts[1:]
                # Process indent syntax if any
                body = []
                for b in bullets:
                    body.append(b)
                data["slides"].append({
                    "type": "standard",
                    "title": slide_title,
                    "body": body
                })
    else:
        # Read from stdin
        try:
            print("No arguments provided. Reading JSON from stdin...", file=sys.stderr)
            data = json.load(sys.stdin)
        except Exception as e:
            print("Error: Must provide either --input, --title/--slides, or pipe JSON via stdin.", file=sys.stderr)
            sys.exit(1)
            
    generate_presentation(data)

if __name__ == "__main__":
    main()
