import sys
from pptx import Presentation

def inspect_layouts(pptx_path):
    prs = Presentation(pptx_path)
    print(f"File: {pptx_path}")
    width = getattr(prs, 'slide_width', 0) or 0
    height = getattr(prs, 'slide_height', 0) or 0
    print(f"Slide Width: {width / 914400} inches")
    print(f"Slide Height: {height / 914400} inches")
    print("\nAvailable Slide Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout Index {i}: {layout.name}")
        # Let's inspect placeholders in each layout
        placeholders = []
        for shape in layout.placeholders:
            placeholders.append(f"{shape.name} (type: {shape.placeholder_format.type}, idx: {shape.placeholder_format.idx})")
        if placeholders:
            print(f"  Placeholders: {', '.join(placeholders)}")
        else:
            print("  No placeholders")

if __name__ == "__main__":
    path = "/Users/mojoaar/AI/skills/new-jysk-pptx/assets/jysk_indoor_fy26.pptx"
    if len(sys.argv) > 1:
        path = sys.argv[1]
    inspect_layouts(path)
